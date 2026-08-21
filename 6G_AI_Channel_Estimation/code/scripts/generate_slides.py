#!/usr/bin/env python3
"""End-to-end implementation slides with figures."""

from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from src.common.utils import load_json

OUT = ROOT / "docs" / "6G_AI_CE_E2E_Implementation.pptx"
PLOTS = ROOT / "outputs" / "plots"
NAVY = RGBColor(0x0B, 0x1C, 0x3D)
CYAN = RGBColor(0x00, 0xC9, 0xFF)
GOLD = RGBColor(0xF5, 0xC4, 0x51)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, l, t, w, h, text, size=18, bold=False, color=WHITE):
    shape = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    return shape


def _add_picture(slide, name, l, t, w):
    path = PLOTS / name
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(l), Inches(t), width=Inches(w))


def main():
    report = {}
    rp = ROOT / "outputs" / "reports" / "train_val_test_report.json"
    if rp.exists():
        report = load_json(rp)
    arch = report.get("architecture", {})

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.25), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()
    _box(s, 0.6, 1.8, 12, 1, "6G AI-Based Secure Channel Estimation", 36, True, CYAN)
    _box(s, 0.6, 2.7, 12, 1, "End-to-end implementation — dataset, multi-agent PHY intelligence,\ndigital twin, attack mitigation, dashboard, and 3GPP-aligned evaluation", 18, False, WHITE)
    _box(s, 0.6, 5.8, 12, 0.8, "3GPP TR 38.901 · TS 38.211 · TR 38.843 · TR 38.811 · Nokia RAN1 R1-2506757 · O-RAN RIC", 14, False, GOLD)
    _box(s, 0.6, 6.5, 12, 0.4, "StrongHER  |  Swetha Kerahalli  |  System Architect, MI RRD AS Algo Innov", 12, False, WHITE)

    # 2 agenda
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _box(s, 0.5, 0.3, 12, 0.6, "Agenda", 28, True, CYAN)
    items = [
        "Problem and 3GPP-aligned opportunity",
        "Architecture: ten agents + digital twin + O-RAN mapping",
        "Synthetic dataset (≥60k) from CDL/TDL, NTN, THz, RIS",
        "Train / validation / test of every model",
        "LS vs MMSE vs AI ensemble — NMSE, BER, SE",
        "Security classification and autonomous mitigation",
        "Digital twin visualizations and closed loop",
        "Dashboard, chatbot, knowledge sources, next steps",
    ]
    _box(s, 0.7, 1.2, 11, 5.5, "\n".join(f"{i+1}.  {t}" for i, t in enumerate(items)), 18)

    # 3 problem
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _box(s, 0.5, 0.3, 12, 0.6, "The 6G CSI problem", 28, True, CYAN)
    _box(s, 0.6, 1.2, 12, 5, "• LS/MMSE break under THz blockage, RIS, ultra-massive MIMO, NTN Doppler\n• Pilot overhead explodes as Nt × Nr and bandwidth grow\n• No native CSI prediction → reactive, not proactive RAN\n• Pilot contamination, jamming, spoofing, poisoning, adversarial CSI\n• Need an AI-native, self-securing channel intelligence plane\n  aligned with Rel-18/19/20 (TR 38.843) and 6G RAN1 studies", 20)

    # 4 architecture
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _box(s, 0.5, 0.3, 12, 0.6, "Multi-agent architecture", 28, True, CYAN)
    _box(s, 0.6, 1.1, 12, 5.8,
         "Channel Agent     CNN+LSTM+Transformer+GNN ensemble CSI estimate\n"
         "CSI Prediction    Reduce pilots when forecast accuracy is high\n"
         "Security Agent    Multi-class attack detection (RF) + IsolationForest\n"
         "Mitigation Agent  Pilot reassign / beam switch / hop / trust / FL retrain\n"
         "Beam Agent        Spatial beam index from SNR, Doppler, array size\n"
         "Mobility Agent    Predictive handover from velocity + neighbor RSRP\n"
         "Optimization      Spectral efficiency vs pilot overhead\n"
         "Digital Twin      Fidelity-gated policy validation\n"
         "Explainability    Permutation importance for operator trust\n"
         "Orchestrator      Global policy fusion for Near-RT RIC", 16)

    # 5 dataset
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _box(s, 0.5, 0.25, 12, 0.5, "3GPP-aligned synthetic dataset", 28, True, CYAN)
    _box(s, 0.5, 0.85, 12, 0.4, f"{arch.get('n_train', 45500)} train  ·  {arch.get('n_validation', 9750)} validation  ·  {arch.get('n_test', 9750)} test   (≥ 60,000 rows)", 16, False, GOLD)
    _add_picture(s, "hist_dataset_splits.png", 0.4, 1.4, 12.4)

    # 6 nmse
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _box(s, 0.5, 0.2, 12, 0.5, "LS vs MMSE vs AI  —  test NMSE / BER / SE", 24, True, CYAN)
    _box(s, 0.5, 0.7, 12, 0.4, f"NMSE Δ {arch.get('nmse_improvement_pct')}%   BER Δ {arch.get('ber_reduction_pct')}%   SE Δ {arch.get('spectral_efficiency_gain_pct')}%", 16, False, GOLD)
    _add_picture(s, "benchmark_ls_mmse_ai.png", 0.4, 1.2, 12.4)

    # 7 cdf
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _box(s, 0.5, 0.2, 12, 0.5, "Estimator NMSE CDFs and scenario CDFs", 24, True, CYAN)
    _add_picture(s, "cdf_nmse_estimators.png", 0.3, 0.9, 6.2)
    _add_picture(s, "cdf_scenario_kpis.png", 6.6, 0.9, 6.4)

    # 8 heatmaps
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _box(s, 0.5, 0.2, 12, 0.5, "Heatmaps — correlation and scenario × profile NMSE", 22, True, CYAN)
    _add_picture(s, "heatmap_feature_correlation.png", 0.3, 0.85, 6.3)
    _add_picture(s, "heatmap_scenario_profile_nmse.png", 6.7, 0.85, 6.3)

    # 9 train val test
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _box(s, 0.5, 0.2, 12, 0.5, "Train / validation / test by agent", 24, True, CYAN)
    _add_picture(s, "model_train_val_test.png", 0.4, 0.9, 12.4)

    # 10 security
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _box(s, 0.5, 0.2, 12, 0.5, "Security agent — classification, ROC, scatter", 24, True, CYAN)
    _add_picture(s, "classification_confusion_matrix.png", 0.2, 0.85, 6.4)
    _add_picture(s, "classification_roc.png", 6.7, 0.85, 6.2)

    # 11 twin
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _box(s, 0.5, 0.2, 12, 0.5, "Digital twin — cell map and fidelity loop", 24, True, CYAN)
    _add_picture(s, "digital_twin_map.png", 0.3, 0.85, 6.3)
    _add_picture(s, "digital_twin_timeseries.png", 6.7, 0.85, 6.3)

    # 12 architecture score
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _box(s, 0.5, 0.2, 12, 0.5, "Overall architecture evaluation", 24, True, CYAN)
    _add_picture(s, "architecture_scorecard.png", 0.3, 0.8, 6.5)
    _add_picture(s, "architecture_radar.png", 7.0, 0.8, 5.8)

    # 13 knowledge
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _box(s, 0.5, 0.3, 12, 0.6, "Knowledge sources", 28, True, CYAN)
    _box(s, 0.6, 1.1, 12, 5.8,
         "3GPP.org   TR 38.901 CDL/TDL · TS 38.211 DMRS/CSI-RS · TR 38.843 AI/ML AiF · TR 38.811 NTN\n\n"
         "System Insights   CFAM RP003187 DMRS channel estimation (5GMax / 5G_L1_2794)\n\n"
         "SharePoint   R1-2506757 6G AI/ML use cases · PHY AI radio deep dive · RAN1#126 · RAN4#120\n\n"
         "Confluence   EE spaces queried (NRAC, RFSW); no indexed hits on this client\n\n"
         "CCFK   Nokia CSF dashboard components (Autonomous RAN pattern)\n\n"
         "agent-shim   clone blocked on scm.cci.nokia.net auth — local shim adapter provided", 16)

    # 14 run
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _box(s, 0.5, 0.3, 12, 0.6, "Run the platform", 28, True, CYAN)
    _box(s, 0.7, 1.3, 12, 5,
         "cd Swetha_StrongHER_projects/6G_AI_Channel_Estimation/code\n"
         "pip install -r requirements.txt\n"
         "python scripts/run_end_to_end.py\n"
         "python scripts/run_api_server.py\n\n"
         "Dashboard   http://localhost:8090/dashboard\n"
         "API         http://localhost:8090/docs\n"
         "Slides      docs/6G_AI_CE_E2E_Implementation.pptx", 20, False, WHITE)

    # 15 close
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _box(s, 0.6, 2.2, 12, 1, "AI-native, self-securing 6G channel intelligence", 28, True, CYAN)
    _box(s, 0.6, 3.3, 12, 1.2, "Accurate CSI  ·  Predicted CSI  ·  Detected attacks  ·  Mitigated threats\nTwin-validated policies  ·  O-RAN ready", 18, False, WHITE)
    _box(s, 0.6, 6.3, 12, 0.4, "Nokia confidential — StrongHER internal use", 12, False, GOLD)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
